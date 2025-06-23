import os
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinterdnd2 import DND_FILES, TkinterDnD
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.dml.color import RGBColor as PPTX_RGBColor
from docx import Document
from docx.shared import RGBColor
import fitz  # PyMuPDF
import math

# ----------------- Helpers -------------------

def create_output_filename(input_path):
    base, ext = os.path.splitext(input_path)
    return f"{base}-Accessible-Copy{ext}"
    
def snap_rect_expand_and_clamp(rect, page_rect, grid_size=10):
    snapped = fitz.Rect(
        math.floor(rect.x0 / grid_size) * grid_size,
        math.floor(rect.y0 / grid_size) * grid_size,
        math.ceil(rect.x1 / grid_size) * grid_size,
        math.ceil(rect.y1 / grid_size) * grid_size,
    )

    # Clamp to the page boundaries
    clamped = fitz.Rect(
        max(snapped.x0, page_rect.x0),
        max(snapped.y0, page_rect.y0),
        min(snapped.x1, page_rect.x1),
        min(snapped.y1, page_rect.y1),
    )
    return clamped
    
def merge_text_blocks(blocks, vertical_threshold=15):
    merged = []
    blocks = sorted([b for b in blocks if b.get("type") == 0], key=lambda b: b["bbox"][1])

    current_block = None
    for block in blocks:
        bbox = block["bbox"]
        rect = fitz.Rect(bbox)

        if current_block is None:
            current_block = block
            continue

        current_bbox = current_block["bbox"]
        current_rect = fitz.Rect(current_bbox)

        # Check vertical proximity (block is right below the current)
        vertical_gap = rect.y0 - current_rect.y1

        if vertical_gap <= vertical_threshold and abs(rect.x0 - current_rect.x0) < 5:
            # Merge: Expand the bbox and combine lines
            current_block["bbox"] = [
                min(current_rect.x0, rect.x0),
                min(current_rect.y0, rect.y0),
                max(current_rect.x1, rect.x1),
                max(current_rect.y1, rect.y1)
            ]
            current_block["lines"] += block.get("lines", [])
        else:
            merged.append(current_block)
            current_block = block

    if current_block:
        merged.append(current_block)

    return merged

def extract_images_with_bboxes(page):
    images = []
    page_rect = page.rect
    for img in page.get_images(full=True):
        xref = img[0]
        try:
            bbox = page.get_image_bbox(xref)
            # Minimum size filter
            if bbox.width < 10 or bbox.height < 10:
                continue

            # Filter images that are almost full page width or height AND aligned to page edges (likely footer/header/sidebar/background)
            near_full_width = bbox.width > page_rect.width * 0.98
            near_full_height = bbox.height > page_rect.height * 0.98
            aligned_left = abs(bbox.x0 - page_rect.x0) < 1
            aligned_right = abs(bbox.x1 - page_rect.x1) < 1
            aligned_top = abs(bbox.y0 - page_rect.y0) < 1
            aligned_bottom = abs(bbox.y1 - page_rect.y1) < 1

            # Exclude if:
            # - full width and aligned top or bottom (header/footer)
            # - full height and aligned left or right (sidebars)
            # - full page (all sides)
            if (
                (near_full_width and (aligned_top or aligned_bottom))
                or (near_full_height and (aligned_left or aligned_right))
                or (near_full_width and near_full_height and aligned_left and aligned_top)
            ):
                continue

            images.append({"xref": xref, "bbox": bbox})
        except Exception:
            continue
    return images


# --------------- DOCX processing -----------------

def process_docx(input_path):
    doc = Document(input_path)

    for para in doc.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            run.font.name = 'Arial'
            run.font.highlight_color = None  # Remove highlighting

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.name = 'Arial'
                        run.font.highlight_color = None

    output_path = create_output_filename(input_path)
    doc.save(output_path)
    print(f"Processed and saved DOCX: {output_path}")

# --------------- PPTX processing -----------------

def safe_set_shape_fill_white(shape):
    fill = shape.fill
    if fill is None:
        return

    fill_type = fill.type

    if fill_type == MSO_FILL.SOLID:
        try:
            transparency = getattr(fill.fore_color, 'transparency', None)
        except Exception:
            fill.solid()
            fill.fore_color.rgb = PPTX_RGBColor(255, 255, 255)
            return

        if transparency is None or transparency < 1.0:
            fill.fore_color.rgb = PPTX_RGBColor(255, 255, 255)

    elif fill_type in (MSO_FILL.PATTERNED, MSO_FILL.GRADIENT, MSO_FILL.PICTURE, MSO_FILL.TEXTURED):
        fill.solid()
        fill.fore_color.rgb = PPTX_RGBColor(255, 255, 255)

    elif fill_type is None:
        return

def process_pptx(input_path):
    prs = Presentation(input_path)

    for slide in prs.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PPTX_RGBColor(255, 255, 255)

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                safe_set_shape_fill_white(shape)

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = PPTX_RGBColor(0, 0, 0)

    output_path = create_output_filename(input_path)
    prs.save(output_path)
    print(f"Processed and saved PPTX: {output_path}")

# --------------- PDF processing -----------------

def process_pdf(input_path, grid_size=10, merge_threshold=15):
    doc = fitz.open(input_path)
    margin = 3  # margin for inset and text shift

    for page_index in range(len(doc)):
        page = doc[page_index]
        rect = page.rect

        # Step 1: Erase existing content by filling white background
        page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1))

        # Step 2: Extract layout structure
        text_dict = page.get_text("dict")
        blocks = text_dict["blocks"]
        images = page.get_images(full=True)

        # Step 3: Insert images first (bottom layer)
        for block in blocks:
            if block.get("type") != 1 or "bbox" not in block:
                continue
            r = fitz.Rect(block["bbox"])

            # Draw black border around images
            page.draw_rect(r, color=(0.99, 0.99, 0.99), width=1)

            # Use the first matching image (buggy but working for now)
            for img in images:
                xref = img[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    page.insert_image(r, pixmap=pix, keep_proportion=True)
                    break
                except Exception as e:
                    print(f"Error inserting image: {e}")
                    continue

        # Step 4: Redraw text blocks and outlines
        for block in blocks:
            if block.get("type") != 0 or "bbox" not in block:
                continue

            r = fitz.Rect(block["bbox"])

            # Skip very small blocks (e.g., artifacts or icons)
            if r.height < 5 or r.width < 10:
                continue

            # Inset the rectangle manually to avoid touching text
            inset_r = fitz.Rect(
                r.x0 + margin,
                r.y0 + margin,
                r.x1 - margin,
                r.y1 - margin
            )

            # Draw light grey outline before text
            page.draw_rect(inset_r, color=(0.9, 0.9, 0.9), width=0.5)

            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if not text:
                        continue
                    x, y = span.get("origin")
                    x += margin
                    y += margin
                    font_size = span.get("size", 10)

                    # Draw white background behind text
                    bg_rect = fitz.Rect(
                        span["bbox"][0] + margin,
                        span["bbox"][1] + margin,
                        span["bbox"][2] + margin,
                        span["bbox"][3] + margin,
                    )
                    page.draw_rect(bg_rect, color=None, fill=(1, 1, 1), overlay=True)

                    # Draw black text on top
                    page.insert_text(
                        (x, y),
                        text,
                        fontname="helv",
                        fontsize=font_size,
                        color=(0, 0, 0),
                    )

    output_path = input_path.replace(".pdf", "-Accessible-Copy.pdf")
    doc.save(output_path)
    print(f"Processed and saved: {output_path}")



# --------------- File handling -----------------

def handle_file(file_path, grid_size=10, merge_threshold=15):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".docx":
            process_docx(file_path)
        elif ext == ".pptx":
            process_pptx(file_path)
        elif ext == ".pdf":
            process_pdf(file_path, grid_size=grid_size, merge_threshold=merge_threshold)
        else:
            messagebox.showerror("Unsupported file", f"File type {ext} is not supported.")
    except Exception as e:
        messagebox.showerror("Error processing file", str(e))


# --------------- GUI -----------------

class AccessibleApp(TkinterDnD.Tk):
    def __init__(self):
        super().__init__()

        self.title("Accessible Document Converter")
        self.geometry("400x420")

        self.label = tk.Label(self, text="Drag & Drop file here\nor use Browse button",
                              width=40, height=10, relief="ridge", borderwidth=2)
        self.label.pack(pady=10)

        self.label.drop_target_register(DND_FILES)
        self.label.dnd_bind('<<Drop>>', self.drop)

        browse_btn = tk.Button(self, text="Browse File", command=self.browse_file)
        browse_btn.pack(pady=10)

        # Slider for Grid Size
        tk.Label(self, text="PDF Grid Size (px)").pack()
        self.grid_size_slider = tk.Scale(self, from_=5, to=50, orient='horizontal')
        self.grid_size_slider.set(20)  # default value
        self.grid_size_slider.pack(fill='x', padx=20)

        # Slider for Merge Sensitivity (vertical threshold)
        tk.Label(self, text="PDF Merge Sensitivity (px)").pack()
        self.merge_sensitivity_slider = tk.Scale(self, from_=5, to=50, orient='horizontal')
        self.merge_sensitivity_slider.set(30)  # default value
        self.merge_sensitivity_slider.pack(fill='x', padx=20)

    def drop(self, event):
        files = self.tk.splitlist(event.data)
        if files:
            file_path = files[0]
            handle_file(
                file_path,
                grid_size=self.grid_size_slider.get(),
                merge_threshold=self.merge_sensitivity_slider.get()
            )

    def browse_file(self):
        file_path = filedialog.askopenfilename(
            filetypes=[("Supported files", "*.pdf *.docx *.pptx")]
        )
        if file_path:
            handle_file(
                file_path,
                grid_size=self.grid_size_slider.get(),
                merge_threshold=self.merge_sensitivity_slider.get()
            )


# --------------- Run -----------------

if __name__ == "__main__":
    try:
        import tkinterdnd2
    except ImportError:
        messagebox.showerror("Missing Dependency",
            "Please install tkinterdnd2 package:\n\npip install tkinterdnd2")
        exit(1)

    app = AccessibleApp()
    app.mainloop()
